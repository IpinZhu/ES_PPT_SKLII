import sys
import json
import os
import win32com.client

def duplicate_slides_win32(pptx_path, content_count):
    # PPT uses 1-based indexing
    # Slide 1: Cover
    # Slide 2: TOC
    # Slide 3: Content 1
    # Slide 4: Content 2
    # Slide 5: Ending
    abs_path = os.path.abspath(pptx_path)
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    presentation = powerpoint.Presentations.Open(abs_path, WithWindow=False)
    
    current_content_slides = 2
    
    try:
        if content_count == 1:
            presentation.Slides(4).Delete()
        elif content_count > 2:
            needed = content_count - 2
            for _ in range(needed):
                presentation.Slides(4).Duplicate()
                
        presentation.Save()
    except Exception as e:
        print(f"Error duplicating slides: {e}")
    finally:
        presentation.Close()
        # powerpoint.Quit() # Keep it open for performance if running multiple times, or quit. We will quit.
        # Check if other presentations are open before quitting
        if powerpoint.Presentations.Count == 0:
            powerpoint.Quit()

def replace_text_in_shape(shape, old_text, new_text):
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if old_text in r.text:
                    r.text = r.text.replace(old_text, str(new_text))

def replace_text_exact(shape, exact_text, new_text):
    if not shape.has_text_frame:
        return False
    # If the entire shape text is exactly `exact_text`
    full_text = "".join([r.text for p in shape.text_frame.paragraphs for r in p.runs]).strip()
    if full_text == exact_text:
        # replace the first run, clear the rest to keep styling
        first = True
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if first:
                    r.text = str(new_text)
                    first = False
                else:
                    r.text = ""
        return True
    return False

def replace_all_text(shape, new_text):
    if not shape.has_text_frame: return
    first = True
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if first:
                r.text = str(new_text)
                first = False
            else:
                r.text = ""

def get_all_shapes(shapes):
    all_shapes = []
    for shape in shapes:
        if shape.shape_type == 6: # GROUP
            all_shapes.extend(get_all_shapes(shape.shapes))
        else:
            all_shapes.append(shape)
    return all_shapes

def process_presentation(json_path, template_path, output_path):
    import shutil
    shutil.copyfile(template_path, output_path)
    
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
        
    slides_data = data.get('slides', [])
    content_count = len(slides_data)
    if content_count == 0:
        print("No content slides to generate.")
        return
        
    duplicate_slides_win32(output_path, content_count)
    
    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation(output_path)
    
    # Slide 1: Cover
    cover_slide = prs.slides[0]
    cover_data = data.get('cover', {})
    for shape in cover_slide.shapes:
        replace_text_in_shape(shape, "大大大大标题", cover_data.get('title', ''))
        replace_text_in_shape(shape, "子标题", cover_data.get('subtitle', ''))
        replace_text_in_shape(shape, "2025年5月", cover_data.get('date', ''))
        if shape.has_text_frame and "申 请 人 ：" in shape.text:
            info = f"申 请 人 ：{cover_data.get('reporter', '')}\n" \
                   f"硕士导师：{cover_data.get('instructor', '')}\n" \
                   f"硕士专业：{cover_data.get('major', '')}"
            replace_all_text(shape, info)
            
    # Slide 2: TOC
    toc_slide = prs.slides[1]
    toc_data = data.get('toc', [])
    all_toc_shapes = get_all_shapes(toc_slide.shapes)
    num_idx = 0
    title_idx = 0
    for shape in all_toc_shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text == "01" and num_idx < len(toc_data):
                replace_all_text(shape, f"{num_idx+1:02d}")
                num_idx += 1
            elif text == "项目背景介绍" and title_idx < len(toc_data):
                replace_all_text(shape, toc_data[title_idx])
                title_idx += 1
                
    # Slide 3 to N+2: Content
    for i, s_data in enumerate(slides_data):
        slide_idx = i + 2
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            replace_text_in_shape(shape, "➢ 核心痛点分析", s_data.get('subtitle', ''))
            
            # Use replace_text_exact or startswith for "01 项目背景介绍"
            if shape.has_text_frame:
                if "项目背景介绍" in shape.text:
                    replace_all_text(shape, s_data.get('title', ''))
                elif "现有方案的局限性较高" in shape.text:
                    bullets = s_data.get('bullets', [])
                    replace_all_text(shape, "\n".join(bullets))
                    
        # Handle image
        if s_data.get('image') and os.path.exists(s_data['image']):
            # insert picture in the middle roughly
            slide.shapes.add_picture(s_data['image'], Inches(4), Inches(2.5), width=Inches(5))
            
    prs.save(output_path)
    print(f"Successfully generated {output_path}")

if __name__ == '__main__':
    if len(sys.argv) < 4:
        print("Usage: python build_from_template.py <data.json> <template.pptx> <output.pptx>")
        sys.exit(1)
    process_presentation(sys.argv[1], sys.argv[2], sys.argv[3])

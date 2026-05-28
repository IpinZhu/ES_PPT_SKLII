import sys
import re
import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def px(pixels):
    """Convert pixels to EMUs (Marp default 1280x720 at 96 DPI)"""
    return int(pixels * 914400 / 96)

def add_title_slide(prs, slide_md):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    
    # Left banner (blue)
    left_banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(100), px(0), px(295), px(720))
    left_banner.fill.solid()
    left_banner.fill.fore_color.rgb = RGBColor(0x00, 0x4E, 0xA1)
    left_banner.line.fill.background()

    # Logo placeholder (we can't easily add the image if it doesn't exist, but we can try)
    logo_path = os.path.join(os.path.dirname(__file__), '..', 'assets', 'heu_logo_badge.png')
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, px(128), px(220), px(240), px(240))

    # Parse title
    titles = re.findall(r'<div class="cover-title-p1">(.*?)</div>', slide_md)
    title_text = titles[0] if len(titles) > 0 else "主标题"
    subtitle_text = titles[1] if len(titles) > 1 else "副标题"
    
    tb = slide.shapes.add_textbox(px(395), px(60), px(885), px(180))
    tf = tb.text_frame
    tf.text = title_text
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if subtitle_text:
        p = tf.add_paragraph()
        p.text = subtitle_text
        p.font.size = Pt(54)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # Meta line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(700), px(340), px(3), px(120))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x00, 0x4E, 0xA1)
    line.line.fill.background()

    # Parse meta
    meta_items = re.findall(r'<span class="label".*?>(.*?)</span>：(.*?)</div>', slide_md)
    meta_y = 340
    for label, value in meta_items:
        tb = slide.shapes.add_textbox(px(730), px(meta_y), px(400), px(40))
        tb.text_frame.text = f"{label}：{value}"
        tb.text_frame.paragraphs[0].font.size = Pt(28)
        meta_y += 50

    # Parse date
    date_match = re.search(r'<div class="cover-date">(.*?)</div>', slide_md)
    date_text = date_match.group(1) if date_match else "2026年X月"
    tb = slide.shapes.add_textbox(px(395), px(575), px(885), px(50))
    tb.text_frame.text = date_text
    tb.text_frame.paragraphs[0].font.size = Pt(28)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_toc_slide(prs, slide_md):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # TOC Title
    tb = slide.shapes.add_textbox(px(100), px(100), px(300), px(200))
    tb.text_frame.text = "目录\nCONTENTS"
    tb.text_frame.paragraphs[0].font.size = Pt(54)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x4E, 0xA1)
    if len(tb.text_frame.paragraphs) > 1:
        tb.text_frame.paragraphs[1].font.size = Pt(24)

    # Parse items
    items = re.findall(r'<span class="text">(.*?)</span>', slide_md)
    y = 150
    for i, item in enumerate(items):
        tb = slide.shapes.add_textbox(px(500), px(y), px(600), px(50))
        tb.text_frame.text = f"{i+1:02d} {item}"
        tb.text_frame.paragraphs[0].font.size = Pt(32)
        y += 70

def add_thanks_slide(prs, slide_md):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Center text
    tb = slide.shapes.add_textbox(px(0), px(250), px(1280), px(200))
    tf = tb.text_frame
    tf.text = "THANK YOU"
    tf.paragraphs[0].font.size = Pt(76)
    tf.paragraphs[0].font.color.rgb = RGBColor(230, 240, 250) # Light watermark
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "谢谢大家"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x4E, 0xA1)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, slide_md):
    # Try to extract heading 1
    h1_match = re.search(r'^#\s+(.*?)$', slide_md, re.MULTILINE)
    title = h1_match.group(1) if h1_match else ""
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only layout
    if slide.shapes.title:
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x4E, 0xA1)
    
    # Strip HTML and get plain text bullets
    text_content = re.sub(r'<[^>]+>', '', slide_md)
    bullets = re.findall(r'^[-\*]\s+(.*?)$', text_content, re.MULTILINE)
    
    if bullets:
        tb = slide.shapes.add_textbox(px(100), px(200), px(1080), px(400))
        tf = tb.text_frame
        for i, b in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = b
            p.font.size = Pt(24)
            p.level = 0

def build_pptx(md_path, pptx_path):
    prs = Presentation()
    prs.slide_width = px(1280)
    prs.slide_height = px(720)
    
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    slides_md = re.split(r'\n---\n+', content)

    for slide_md in slides_md:
        slide_md = slide_md.strip()
        if not slide_md or slide_md.startswith('marp: true'):
            continue
            
        if 'class: title-slide' in slide_md:
            add_title_slide(prs, slide_md)
        elif 'class: toc-slide' in slide_md:
            add_toc_slide(prs, slide_md)
        elif 'class: thanks-slide' in slide_md:
            add_thanks_slide(prs, slide_md)
        else:
            add_content_slide(prs, slide_md)

    prs.save(pptx_path)
    print(f"Editable PPTX saved to {pptx_path}")

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python build_editable_pptx.py <input.md> [output.pptx]")
        sys.exit(1)
    input_md = sys.argv[1]
    output_pptx = sys.argv[2] if len(sys.argv) > 2 else input_md.replace('.md', '.pptx')
    build_pptx(input_md, output_pptx)

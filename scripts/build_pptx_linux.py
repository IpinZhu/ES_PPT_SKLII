"""Linux-friendly PPTX generator for the es-ppt-skill template (no PowerPoint required).

Phase 2 engine. Produces a fully editable .pptx by:
  * cloning slide 3 (normal-layout) and slide 4 (highlight-layout) once per JSON entry
  * reordering so layout becomes [cover, toc, ...content..., thanks]
  * filling every text frame's first paragraph in place (preserving pPr / first-run rPr)

Math support
------------
Every text we inject is scanned for inline LaTeX delimited by ``$ ... $``.
Each math segment is converted with latex2mathml -> mathml2omml and embedded
into the paragraph as ``<a14:m><m:oMath>...</m:oMath></a14:m>`` so PowerPoint /
WPS / Keynote render it as a native equation, not as raw ``$...$`` text.
Plain text segments keep using ``<a:r>`` runs and inherit the original rPr.
"""
from __future__ import annotations
import sys
import os
import json
import shutil
import copy
import re
from pptx import Presentation
from pptx.oxml.ns import qn, nsmap
from pptx.opc.package import PackURI
from pptx.parts.slide import SlidePart
from lxml import etree


PRES_SLIDE_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"

# Namespaces used for embedding OMML inside DrawingML paragraphs.
A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"


def _qn(prefix: str, tag: str, ns_uri: str) -> str:
    return f"{{{ns_uri}}}{tag}"


# ---------- LaTeX -> OMML ----------
try:
    import latex2mathml.converter as _l2m
    import mathml2omml as _m2o
    HAS_MATH = True
except Exception:  # pragma: no cover
    HAS_MATH = False


_MATH_RE = re.compile(r"\$([^$\n]+)\$")


def latex_to_omml_element(latex_src: str):
    """Return an lxml Element for ``<m:oMath>`` representing *latex_src*.

    Returns None if conversion fails or math libs are missing — callers must
    fall back to plain text in that case.
    """
    if not HAS_MATH:
        return None
    try:
        mml = _l2m.convert(latex_src)
        omml = _m2o.convert(mml)
    except Exception:
        return None
    # mathml2omml returns a string starting with <m:oMath>...; wrap it with a
    # stub root that declares the m: namespace so lxml can parse it.
    wrapped = (
        f'<root xmlns:m="{M_NS}">{omml}</root>'
    )
    try:
        root = etree.fromstring(wrapped.encode("utf-8"))
    except etree.XMLSyntaxError:
        return None
    omath_nodes = root.findall(_qn("m", "oMath", M_NS))
    if not omath_nodes:
        return None
    return omath_nodes[0]


def split_text_with_math(text: str):
    """Yield ``("text", str)`` and ``("math", latex_src)`` segments."""
    if not text:
        return
    pos = 0
    for m in _MATH_RE.finditer(text):
        if m.start() > pos:
            yield "text", text[pos:m.start()]
        yield "math", m.group(1)
        pos = m.end()
    if pos < len(text):
        yield "text", text[pos:]


# ---------- slide cloning / reordering ----------
def clone_slide(prs, source_index):
    src_part = prs.slides[source_index].part
    pkg = prs.part.package
    used = {p.partname.lower() for p in pkg.iter_parts()}
    i = 1
    while True:
        candidate = PackURI(f"/ppt/slides/slide{i}.xml")
        if str(candidate).lower() not in used:
            break
        i += 1
    new_element = copy.deepcopy(src_part._element)
    new_part = SlidePart(candidate, src_part.content_type, pkg, new_element)
    for rel in src_part.rels.values():
        if rel.is_external:
            new_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_part.rels.get_or_add(rel.reltype, rel.target_part)
    rId = prs.part.relate_to(new_part, PRES_SLIDE_RELTYPE)
    sldIdLst = prs.slides._sldIdLst
    existing_ids = [int(e.get("id")) for e in sldIdLst.findall(qn("p:sldId"))]
    new_id = max(existing_ids + [255]) + 1
    sldId = etree.SubElement(sldIdLst, qn("p:sldId"))
    sldId.set("id", str(new_id))
    sldId.set(qn("r:id"), rId)
    if hasattr(prs.slides, "_slides"):
        prs.slides.__dict__.pop("_slides", None)
    return prs.slides[len(prs.slides) - 1]


def remove_slide(prs, slide_index):
    sldIdLst = prs.slides._sldIdLst
    entries = list(sldIdLst.findall(qn("p:sldId")))
    target = entries[slide_index]
    rId = target.get(qn("r:id"))
    sldIdLst.remove(target)
    prs.part.drop_rel(rId)
    if hasattr(prs.slides, "_slides"):
        prs.slides.__dict__.pop("_slides", None)


def move_slide(prs, src_index, dst_index):
    sldIdLst = prs.slides._sldIdLst
    entries = list(sldIdLst.findall(qn("p:sldId")))
    target = entries[src_index]
    sldIdLst.remove(target)
    entries = list(sldIdLst.findall(qn("p:sldId")))
    if dst_index >= len(entries):
        sldIdLst.append(target)
    else:
        sldIdLst.insert(dst_index, target)
    if hasattr(prs.slides, "_slides"):
        prs.slides.__dict__.pop("_slides", None)


# ---------- text helpers ----------
def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == 6:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def _capture_first_run_rPr(text_frame):
    paras = list(text_frame.paragraphs)
    if not paras:
        return None, None
    first = paras[0]
    pPr = first._p.find(qn("a:pPr"))
    pPr_copy = copy.deepcopy(pPr) if pPr is not None else None
    rPr_copy = None
    for r in first.runs:
        rPr_node = r._r.find(qn("a:rPr"))
        if rPr_node is not None:
            rPr_copy = copy.deepcopy(rPr_node)
            break
    return pPr_copy, rPr_copy


def _add_run_with_text(p_elem, text, rPr_copy):
    r = etree.SubElement(p_elem, qn("a:r"))
    if rPr_copy is not None:
        r.append(copy.deepcopy(rPr_copy))
    t = etree.SubElement(r, qn("a:t"))
    t.text = text


def _add_math(p_elem, latex_src):
    """Append <a14:m><m:oMath/></a14:m> to *p_elem*. Returns True on success."""
    omath = latex_to_omml_element(latex_src)
    if omath is None:
        return False
    a14_m = etree.SubElement(p_elem, _qn("a14", "m", A14_NS), nsmap={"a14": A14_NS, "m": M_NS})
    a14_m.append(omath)
    return True


def _fill_paragraph(p_elem, text, rPr_copy):
    """Populate paragraph *p_elem* with mixed text + math runs based on *text*.

    A leading <a:pPr> child, if any, is preserved.
    Existing <a:r>, <a:br>, <a:fld>, <a14:m> children are removed first.
    """
    for child in list(p_elem):
        if child.tag in (qn("a:r"), qn("a:br"), qn("a:fld"), _qn("a14", "m", A14_NS)):
            p_elem.remove(child)
    for kind, payload in split_text_with_math(text):
        if kind == "text":
            if payload:
                _add_run_with_text(p_elem, payload, rPr_copy)
        else:  # math
            if not _add_math(p_elem, payload):
                _add_run_with_text(p_elem, f"${payload}$", rPr_copy)


def set_text_keep_format(shape, new_text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    pPr_copy, rPr_copy = _capture_first_run_rPr(tf)
    txBody = tf._txBody
    for child in list(txBody):
        if child.tag in (qn("a:bodyPr"), qn("a:lstStyle")):
            continue
        txBody.remove(child)
    p = etree.SubElement(txBody, qn("a:p"))
    if pPr_copy is not None:
        p.append(copy.deepcopy(pPr_copy))
    _fill_paragraph(p, new_text, rPr_copy)


def set_multiline_text(shape, lines, line_height_emu=457200, min_lines=3):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    pPr_copy, rPr_copy = _capture_first_run_rPr(tf)
    txBody = tf._txBody
    for child in list(txBody):
        if child.tag in (qn("a:bodyPr"), qn("a:lstStyle")):
            continue
        txBody.remove(child)
    for line in lines:
        p = etree.SubElement(txBody, qn("a:p"))
        if pPr_copy is not None:
            p.append(copy.deepcopy(pPr_copy))
        _fill_paragraph(p, line, rPr_copy)
    # auto-fit shape height to number of lines (only grow, never shrink)
    n = max(len(lines), min_lines)
    desired = n * line_height_emu
    # modify XML directly for reliable shape height change
    # find <p:spPr><a:xfrm><a:ext cy="..."/>
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is not None:
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is not None:
            ext = xfrm.find(qn("a:ext"))
            if ext is not None:
                try:
                    current_cy = int(ext.get("cy", 0))
                    if desired > current_cy:
                        ext.set("cy", str(desired))
                except (ValueError, TypeError):
                    pass
    # disable PowerPoint auto-shrink so text stays at original size
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        # remove any existing autofit children
        for tag in ("a:noAutofit", "a:normAutofit", "a:spAutoFit"):
            for n_ in bodyPr.findall(qn(tag)):
                bodyPr.remove(n_)
        etree.SubElement(bodyPr, qn("a:noAutofit"))
        bodyPr.set("wrap", "square")


def set_two_line_with_break(shape, line1, line2):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    pPr_copy, rPr_copy = _capture_first_run_rPr(tf)
    txBody = tf._txBody
    for child in list(txBody):
        if child.tag in (qn("a:bodyPr"), qn("a:lstStyle")):
            continue
        txBody.remove(child)
    p = etree.SubElement(txBody, qn("a:p"))
    if pPr_copy is not None:
        p.append(copy.deepcopy(pPr_copy))
    _fill_paragraph(p, line1, rPr_copy)
    etree.SubElement(p, qn("a:br"))
    # use a separate temp paragraph to render line2 into the same paragraph
    tmp = etree.SubElement(txBody, qn("a:p"))
    _fill_paragraph(tmp, line2, rPr_copy)
    # move all of tmp's children (except pPr) into p, then drop tmp
    for child in list(tmp):
        if child.tag == qn("a:pPr"):
            continue
        p.append(child)
    txBody.remove(tmp)


# ---------- fillers ----------
def fill_cover(slide, cover):
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if "大大大大标题" in text or "子标题" in text:
            set_two_line_with_break(shape, cover.get("title", ""), cover.get("subtitle", ""))
        elif "申 请 人 ：" in text:
            lines = [
                f"汇 报 人 ：{cover.get('reporter','')}",
                f"指导老师：{cover.get('instructor','')}",
                f"专 业 ：{cover.get('major','')}",
                f"日 期 ：{cover.get('date','')}",
            ]
            set_multiline_text(shape, lines)
        elif text.strip().endswith("年5月") or text.strip().startswith("2025年") or text.strip().startswith("2026年"):
            set_text_keep_format(shape, cover.get("date", ""))


def fill_toc(slide, toc):
    shapes = list(iter_shapes(slide.shapes))
    num_shapes = []
    title_shapes = []
    for sh in shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t == "01":
            num_shapes.append(sh)
        elif t == "项目背景介绍":
            title_shapes.append(sh)
    n = min(len(toc), len(num_shapes), len(title_shapes))
    for i in range(n):
        set_text_keep_format(num_shapes[i], f"{i+1:02d}")
        set_text_keep_format(title_shapes[i], toc[i])


def fill_content(slide, s_data):
    bullets = s_data.get("bullets", [])
    title = s_data.get("title", "")
    subtitle = s_data.get("subtitle", "")
    highlight_text = s_data.get("highlight")
    layout = s_data.get("layout", "normal")
    for shape in iter_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        stripped = txt.strip()
        if "项目背景介绍" in stripped:
            set_text_keep_format(shape, title)
        elif stripped.startswith("➢"):
            set_text_keep_format(shape, subtitle)
        elif "现有方案的局限性较高" in txt or "用户对高颜值" in txt:
            set_multiline_text(shape, bullets)
        elif layout == "highlight" and "*HIGHLIGHT" in txt:
            set_text_keep_format(shape, highlight_text or (bullets[0] if bullets else ""))




def _hoist_math_namespaces(pptx_path):
    """Move xmlns:a14 / xmlns:m declarations from each <a14:m> up to the root
    <p:sld>, so the file is smaller and matches the layout PowerPoint emits."""
    import zipfile
    import re as _re
    A14 = "http://schemas.microsoft.com/office/drawing/2010/main"
    M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    with zipfile.ZipFile(pptx_path, "r") as zin:
        items = {n: zin.read(n) for n in zin.namelist()}
    changed = False
    for name in list(items):
        if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
            continue
        raw = items[name].decode("utf-8")
        if "<a14:m " not in raw:
            continue
        def add_ns(m):
            head = m.group(1)
            new = head
            if "xmlns:a14=" not in new:
                new += f' xmlns:a14="{A14}"'
            if "xmlns:m=" not in new:
                new += f' xmlns:m="{M}"'
            return new + ">"
        raw = _re.sub(r"(<p:sld\b[^>]*?)>", add_ns, raw, count=1)
        raw = raw.replace(f'<a14:m xmlns:a14="{A14}" xmlns:m="{M}">', "<a14:m>")
        items[name] = raw.encode("utf-8")
        changed = True
    if not changed:
        return
    with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for n, b in items.items():
            zout.writestr(n, b)


def process(json_path, template_path, output_path):
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    shutil.copyfile(template_path, output_path)
    prs = Presentation(output_path)
    NORMAL_IDX = 2
    HIGHLIGHT_IDX = 3
    slides_data = data.get("slides", [])
    for s in slides_data:
        idx = HIGHLIGHT_IDX if s.get("layout") == "highlight" else NORMAL_IDX
        clone_slide(prs, idx)
        last = len(prs.slides) - 1
        move_slide(prs, last, last - 1)
    remove_slide(prs, 2)
    remove_slide(prs, 2)
    fill_cover(prs.slides[0], data.get("cover", {}))
    fill_toc(prs.slides[1], data.get("toc", []))
    for i, s_data in enumerate(slides_data):
        fill_content(prs.slides[2 + i], s_data)
    prs.save(output_path)
    _hoist_math_namespaces(output_path)
    print(f"Successfully generated {output_path}")
    if not HAS_MATH:
        print("Note: latex2mathml / mathml2omml not installed — math fell back to plain text. "
              "Run: pip install latex2mathml mathml2omml")


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python build_pptx_linux.py <data.json> <template.pptx> <output.pptx>")
        sys.exit(1)
    process(sys.argv[1], sys.argv[2], sys.argv[3])
